"""
Módulo para extracción exhaustiva de URLs de archivos PowerPoint (.pptx)
"""

import re
import zipfile
import xml.etree.ElementTree as ET
from pptx import Presentation
from urllib.parse import urlparse
import io

class PPTXURLExtractor:
    """Clase para extraer URLs de manera exhaustiva de archivos PPTX"""
    
    def __init__(self):
        # Expresiones regulares mejoradas para detectar URLs COMPLETAS sin división
        self.url_patterns = [
            # URLs completas con http/https - MEJORADO para capturar URLs completas
            re.compile(r'https?://[^\s<>"\']+', re.IGNORECASE),
            # URLs con www sin protocolo - MEJORADO
            re.compile(r'www\.[^\s<>"\']+', re.IGNORECASE),
            # URLs tipo ftp - MEJORADO
            re.compile(r'ftp://[^\s<>"\']+', re.IGNORECASE),
            # Emails (mantener simple)
            re.compile(r'[a-zA-Z0-9._%+-]+@[a-zA-Z0-9.-]+\.[a-zA-Z]{2,}', re.IGNORECASE),
            # URLs en formato de enlace tipo [texto](url)
            re.compile(r'\[([^\]]+)\]\((https?://[^\)]+)\)', re.IGNORECASE),
            # URLs que empiezan con dominios conocidos sin protocolo - MEJORADO
            re.compile(r'(?:^|\s)([a-zA-Z0-9-]+\.(?:com|org|net|edu|gov|mil|int|co|io|me|ly|tk|cc|tv|fm|am|info|biz|name|pro|museum|aero|coop|jobs|travel|mobi|asia|cat|tel|xxx|post|geo|local|youtube|google|github|microsoft|amazon|facebook|twitter|linkedin|instagram|tiktok|vimeo|zoom|teams)[^\s<>"\']*)', re.IGNORECASE),
            # IPs con puertos y rutas - MEJORADO
            re.compile(r'(?:https?://)?(?:[0-9]{1,3}\.){3}[0-9]{1,3}(?::[0-9]+)?(?:/[^\s<>"\']*)?', re.IGNORECASE)
        ]
        
        # Patrón principal MÁS AMPLIO para capturar URLs completas
        self.url_pattern = re.compile(
            r'(?:https?://|ftp://|www\.)[^\s<>"\']+',
            re.IGNORECASE
        )
        
        # Namespaces XML de PowerPoint
        self.namespaces = {
            'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
            'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
            'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'
        }
    
    def extract_urls_from_file(self, file_path_or_content):
        """
        Extraer URLs de un archivo PPTX
        
        Args:
            file_path_or_content: Ruta del archivo o contenido en bytes
            
        Returns:
            List[dict]: Lista de URLs encontradas con contexto detallado
        """
        urls_found = []
        
        try:
            # Abrir el archivo PPTX
            if isinstance(file_path_or_content, str):
                prs = Presentation(file_path_or_content)
                with open(file_path_or_content, 'rb') as f:
                    zip_content = f.read()
            else:
                prs = Presentation(io.BytesIO(file_path_or_content))
                zip_content = file_path_or_content
            
            # Método 1: Extraer URLs usando python-pptx (texto visible y shapes)
            urls_found.extend(self._extract_from_presentation_object(prs))
            
            # Método 2: Extraer URLs del archivo ZIP/XML (búsqueda exhaustiva)
            urls_found.extend(self._extract_from_xml_content(zip_content))
            
            # Método 3: NUEVO - Búsqueda brutal en todo el contenido como último recurso
            urls_found.extend(self._extract_from_all_content_brute_force(zip_content))
            
            # DEDUPLICACIÓN MEJORADA Y ROBUSTA
            unique_urls = self._deduplicate_urls_advanced(urls_found)
            
            return unique_urls
            
        except Exception as e:
            print(f"Error al procesar archivo PPTX: {str(e)}")
            return []
    
    def _extract_from_presentation_object(self, prs):
        """Extraer URLs usando el objeto Presentation de python-pptx con búsqueda exhaustiva en shapes"""
        urls_found = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            # Extraer de formas en la diapositiva
            urls_found.extend(self._extract_from_shapes(slide.shapes, slide_num))
            
            # Notas de la diapositiva
            if slide.has_notes_slide:
                notes_text = slide.notes_slide.notes_text_frame.text
                urls_in_notes = self._find_urls_in_text(notes_text)
                for url in urls_in_notes:
                    urls_found.append({
                        'url': url,
                        'location': f'Diapositiva {slide_num} - Notas',
                        'context': notes_text[:100] + '...' if len(notes_text) > 100 else notes_text
                    })
        
        return urls_found
    
    def _extract_from_shapes(self, shapes, slide_num, parent_context=""):
        """Extraer URLs de shapes de manera recursiva con búsqueda EXHAUSTIVA en todos los elementos"""
        urls_found = []
        
        for shape_idx, shape in enumerate(shapes):
            shape_context = f"{parent_context}Shape {shape_idx + 1}"
            
            try:
                # 1. TEXTO DIRECTO EN FORMAS
                if hasattr(shape, 'text') and shape.text:
                    urls_in_text = self._find_urls_in_text(shape.text)
                    for url in urls_in_text:
                        urls_found.append({
                            'url': url,
                            'location': f'Diapositiva {slide_num} - {shape_context} - Texto directo',
                            'context': shape.text[:100] + '...' if len(shape.text) > 100 else shape.text
                        })
                
                # 2. HIPERVÍNCULOS EN CLICK ACTIONS
                if hasattr(shape, 'click_action'):
                    try:
                        if hasattr(shape.click_action, 'hyperlink') and shape.click_action.hyperlink:
                            hyperlink = shape.click_action.hyperlink
                            if hasattr(hyperlink, 'address') and hyperlink.address:
                                urls_found.append({
                                    'url': hyperlink.address,
                                    'location': f'Diapositiva {slide_num} - {shape_context} - Hipervínculo de acción',
                                    'context': f'Click action hyperlink: {hyperlink.address}'
                                })
                    except Exception as e:
                        pass  # Algunos shapes no tienen click_action válido
                
                # 3. TEXT_FRAME CON BÚSQUEDA EXHAUSTIVA
                if hasattr(shape, 'text_frame') and shape.text_frame:
                    try:
                        # Texto completo del text_frame
                        if hasattr(shape.text_frame, 'text') and shape.text_frame.text:
                            urls_in_frame = self._find_urls_in_text(shape.text_frame.text)
                            for url in urls_in_frame:
                                urls_found.append({
                                    'url': url,
                                    'location': f'Diapositiva {slide_num} - {shape_context} - Text Frame',
                                    'context': shape.text_frame.text[:100] + '...' if len(shape.text_frame.text) > 100 else shape.text_frame.text
                                })
                        
                        # Párrafos individuales
                        for para_idx, paragraph in enumerate(shape.text_frame.paragraphs):
                            if hasattr(paragraph, 'text') and paragraph.text:
                                urls_in_para = self._find_urls_in_text(paragraph.text)
                                for url in urls_in_para:
                                    urls_found.append({
                                        'url': url,
                                        'location': f'Diapositiva {slide_num} - {shape_context} - Párrafo {para_idx + 1}',
                                        'context': paragraph.text[:100] + '...' if len(paragraph.text) > 100 else paragraph.text
                                    })
                            
                            # Runs individuales con hipervínculos
                            for run_idx, run in enumerate(paragraph.runs):
                                # Texto del run
                                if hasattr(run, 'text') and run.text:
                                    urls_in_run = self._find_urls_in_text(run.text)
                                    for url in urls_in_run:
                                        urls_found.append({
                                            'url': url,
                                            'location': f'Diapositiva {slide_num} - {shape_context} - Run {run_idx + 1}',
                                            'context': run.text[:100] + '...' if len(run.text) > 100 else run.text
                                        })
                                
                                # Hipervínculos en runs
                                try:
                                    if hasattr(run, 'hyperlink') and run.hyperlink:
                                        if hasattr(run.hyperlink, 'address') and run.hyperlink.address:
                                            urls_found.append({
                                                'url': run.hyperlink.address,
                                                'location': f'Diapositiva {slide_num} - {shape_context} - Hipervínculo en run {run_idx + 1}',
                                                'context': run.text or 'Hipervínculo sin texto visible'
                                            })
                                except Exception:
                                    pass
                    except Exception:
                        pass
                
                # 4. FORMAS AGRUPADAS (GroupShape) - RECURSIÓN PROFUNDA
                if hasattr(shape, 'shapes'):
                    try:
                        nested_urls = self._extract_from_shapes(shape.shapes, slide_num, f"{shape_context} - Grupo - ")
                        urls_found.extend(nested_urls)
                    except Exception:
                        pass
                
                # 5. TABLAS
                if hasattr(shape, 'table'):
                    try:
                        table_urls = self._extract_from_table(shape.table, slide_num, shape_context)
                        urls_found.extend(table_urls)
                    except Exception:
                        pass
                
                # 6. ATRIBUTOS DEL SHAPE (nombre, alt text, etc.)
                try:
                    # Nombre del shape
                    if hasattr(shape, 'name') and shape.name:
                        urls_in_name = self._find_urls_in_text(shape.name)
                        for url in urls_in_name:
                            urls_found.append({
                                'url': url,
                                'location': f'Diapositiva {slide_num} - {shape_context} - Nombre',
                                'context': f'Nombre del shape: {shape.name}'
                            })
                    
                    # Texto alternativo
                    if hasattr(shape, 'element'):
                        # Buscar en atributos del elemento XML
                        xml_str = ET.tostring(shape.element, encoding='unicode')
                        urls_in_xml = self._find_urls_in_text(xml_str)
                        for url in urls_in_xml:
                            urls_found.append({
                                'url': url,
                                'location': f'Diapositiva {slide_num} - {shape_context} - Atributos XML',
                                'context': 'Encontrado en atributos XML del elemento'
                            })
                except Exception:
                    pass
                
                # 7. SMART ART y DIAGRAMAS
                try:
                    if hasattr(shape, 'element') and shape.element is not None:
                        # Buscar en el XML del elemento para Smart Art
                        element_xml = ET.tostring(shape.element, encoding='unicode')
                        if 'dgm:' in element_xml or 'smartArt' in element_xml:
                            urls_in_smartart = self._find_urls_in_text(element_xml)
                            for url in urls_in_smartart:
                                urls_found.append({
                                    'url': url,
                                    'location': f'Diapositiva {slide_num} - {shape_context} - SmartArt',
                                    'context': 'Encontrado en SmartArt/Diagrama'
                                })
                except Exception:
                    pass
                
                # 8. CONTENIDO MULTIMEDIA (si tiene URLs embebidas)
                try:
                    if hasattr(shape, 'shape_type'):
                        # Videos, audios, objetos embebidos
                        if 'MEDIA' in str(shape.shape_type) or 'OLE' in str(shape.shape_type):
                            # Buscar en el elemento XML para URLs de multimedia
                            if hasattr(shape, 'element'):
                                media_xml = ET.tostring(shape.element, encoding='unicode')
                                urls_in_media = self._find_urls_in_text(media_xml)
                                for url in urls_in_media:
                                    urls_found.append({
                                        'url': url,
                                        'location': f'Diapositiva {slide_num} - {shape_context} - Multimedia',
                                        'context': f'URL en elemento multimedia ({shape.shape_type})'
                                    })
                except Exception:
                    pass
                
                # 9. BUSCAR EN TODAS LAS PROPIEDADES DEL SHAPE
                try:
                    # Iterar sobre todas las propiedades disponibles
                    for attr_name in dir(shape):
                        if not attr_name.startswith('_'):
                            try:
                                attr_value = getattr(shape, attr_name)
                                if isinstance(attr_value, str) and len(attr_value) > 10:
                                    urls_in_attr = self._find_urls_in_text(attr_value)
                                    for url in urls_in_attr:
                                        urls_found.append({
                                            'url': url,
                                            'location': f'Diapositiva {slide_num} - {shape_context} - Propiedad {attr_name}',
                                            'context': f'{attr_name}: {attr_value[:100]}...' if len(attr_value) > 100 else f'{attr_name}: {attr_value}'
                                        })
                            except Exception:
                                pass
                except Exception:
                    pass
                
            except Exception as e:
                # Log del error para debug
                pass
        
        return urls_found
    
    def _extract_from_table(self, table, slide_num, shape_context):
        """Extraer URLs de tablas"""
        urls_found = []
        
        for row_idx, row in enumerate(table.rows):
            for col_idx, cell in enumerate(row.cells):
                if cell.text:
                    urls_in_cell = self._find_urls_in_text(cell.text)
                    for url in urls_in_cell:
                        urls_found.append({
                            'url': url,
                            'location': f'Diapositiva {slide_num} - {shape_context} - Tabla celda ({row_idx + 1},{col_idx + 1})',
                            'context': cell.text[:100] + '...' if len(cell.text) > 100 else cell.text
                        })
                
                # Hipervínculos en celdas de tabla
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if hasattr(run, 'hyperlink') and run.hyperlink and hasattr(run.hyperlink, 'address') and run.hyperlink.address:
                            urls_found.append({
                                'url': run.hyperlink.address,
                                'location': f'Diapositiva {slide_num} - {shape_context} - Tabla celda ({row_idx + 1},{col_idx + 1}) - Hipervínculo',
                                'context': run.text or 'Hipervínculo en tabla'
                            })
        
        return urls_found
    
    def _extract_from_xml_content(self, zip_content):
        """Extraer URLs directamente del contenido XML del archivo PPTX con información de slides"""
        urls_found = []
        
        try:
            with zipfile.ZipFile(io.BytesIO(zip_content), 'r') as zip_file:
                # Buscar en archivos de slides específicos
                for file_name in zip_file.namelist():
                    if file_name.startswith('ppt/slides/slide') and file_name.endswith('.xml'):
                        # Extraer número de slide del nombre del archivo
                        slide_match = re.search(r'slide(\d+)\.xml', file_name)
                        slide_num = int(slide_match.group(1)) if slide_match else 0
                        
                        try:
                            xml_content = zip_file.read(file_name).decode('utf-8')
                            
                            # Buscar URLs en el contenido XML
                            urls_in_xml = self._find_urls_in_text(xml_content)
                            for url in urls_in_xml:
                                urls_found.append({
                                    'url': url,
                                    'location': f'Diapositiva {slide_num} - XML interno',
                                    'context': 'Encontrado en XML de diapositiva'
                                })
                            
                            # Buscar URLs en atributos XML específicos
                            urls_found.extend(self._extract_from_xml_attributes(
                                xml_content, f'Diapositiva {slide_num}', file_name
                            ))
                            
                        except UnicodeDecodeError:
                            continue
                
                # Buscar en archivos de relaciones (_rels) - AQUÍ ES DONDE ESTÁN MUCHOS HIPERVÍNCULOS
                for file_name in zip_file.namelist():
                    if '_rels' in file_name and file_name.endswith('.rels'):
                        try:
                            xml_content = zip_file.read(file_name).decode('utf-8')
                            
                            # Los archivos .rels contienen los hipervínculos externos
                            # Buscar elementos <Relationship> con Type="hyperlink"
                            hyperlink_pattern = r'<Relationship[^>]*Type="[^"]*hyperlink[^"]*"[^>]*Target="([^"]+)"'
                            hyperlink_matches = re.findall(hyperlink_pattern, xml_content, re.IGNORECASE)
                            
                            for target_url in hyperlink_matches:
                                if self._is_valid_url(target_url):
                                    # Determinar a qué slide corresponde este archivo de relación
                                    slide_num = 0
                                    if 'slides/_rels/slide' in file_name:
                                        slide_match = re.search(r'slide(\d+)\.xml\.rels', file_name)
                                        slide_num = int(slide_match.group(1)) if slide_match else 0
                                    
                                    urls_found.append({
                                        'url': target_url,
                                        'location': f'Diapositiva {slide_num} - Hipervínculo' if slide_num > 0 else 'Archivo de relaciones',
                                        'context': f'Hipervínculo externo desde {file_name}'
                                    })
                            
                            # También buscar URLs en cualquier parte del contenido de relaciones
                            urls_in_rels = self._find_urls_in_text(xml_content)
                            for url in urls_in_rels:
                                # Evitar duplicados de los ya encontrados
                                if url not in [h['url'] for h in urls_found]:
                                    slide_num = 0
                                    if 'slides/_rels/slide' in file_name:
                                        slide_match = re.search(r'slide(\d+)\.xml\.rels', file_name)
                                        slide_num = int(slide_match.group(1)) if slide_match else 0
                                    
                                    urls_found.append({
                                        'url': url,
                                        'location': f'Diapositiva {slide_num} - Relación' if slide_num > 0 else 'Archivo de relaciones',
                                        'context': f'Encontrado en relaciones: {file_name}'
                                    })
                            
                        except UnicodeDecodeError:
                            continue
                
                # Buscar en otros archivos XML que pueden contener URLs
                xml_files_to_check = [
                    'ppt/presentation.xml',  # Presentación principal SOLAMENTE
                ]
                
                for xml_file in xml_files_to_check:
                    if xml_file in zip_file.namelist():
                        try:
                            xml_content = zip_file.read(xml_file).decode('utf-8')
                            urls_in_file = self._find_urls_in_text(xml_content)
                            for url in urls_in_file:
                                # VERIFICAR QUE NO SEA METADATA ANTES DE AGREGAR
                                if self._is_valid_url(url):
                                    urls_found.append({
                                        'url': url,
                                        'location': f'Archivo {xml_file}',
                                        'context': f'Encontrado en {xml_file}'
                                    })
                        except UnicodeDecodeError:
                            continue
        
        except Exception as e:
            print(f"Error al procesar contenido XML: {str(e)}")
        
        return urls_found
    
    def _extract_from_xml_attributes(self, xml_content, location_prefix, file_name):
        """Extraer URLs de atributos XML específicos con búsqueda exhaustiva"""
        urls_found = []
        
        # Atributos XML donde pueden estar las URLs
        url_attributes = [
            r'r:id="([^"]*)"',  # IDs de relaciones
            r'action="([^"]*)"',  # Acciones de hipervínculos
            r'r:href="([^"]*)"',  # Enlaces de relaciones
            r'target="([^"]*)"',  # Objetivos de enlaces
            r'tooltip="([^"]*)"',  # Tooltips que pueden contener URLs
            r'description="([^"]*)"',  # Descripciones
            r'title="([^"]*)"',  # Títulos
            r'uri="([^"]*)"',  # URIs
            r'val="([^"]*)"',  # Valores que pueden ser URLs
            r'text="([^"]*)"',  # Texto que puede contener URLs
        ]
        
        for attr_pattern in url_attributes:
            matches = re.findall(attr_pattern, xml_content, re.IGNORECASE)
            for match in matches:
                urls_in_attr = self._find_urls_in_text(match)
                for url in urls_in_attr:
                    urls_found.append({
                        'url': url,
                        'location': f'{location_prefix} - Atributo XML',
                        'context': f'Encontrado en atributo: {match[:100]}...' if len(match) > 100 else f'Encontrado en atributo: {match}'
                    })
        
        # Buscar también en elementos de texto XML
        text_elements = [
            r'<a:t>([^<]+)</a:t>',  # Texto en elementos drawingML
            r'<a:p>([^<]+)</a:p>',  # Párrafos en drawingML
            r'<p:ph>([^<]+)</p:ph>',  # Placeholders
            r'<p:cNvPr[^>]*name="([^"]*)"',  # Nombres de elementos
            r'<p:cNvPr[^>]*descr="([^"]*)"',  # Descripciones de elementos
        ]
        
        for text_pattern in text_elements:
            matches = re.findall(text_pattern, xml_content, re.IGNORECASE | re.DOTALL)
            for match in matches:
                urls_in_text = self._find_urls_in_text(match)
                for url in urls_in_text:
                    urls_found.append({
                        'url': url,
                        'location': f'{location_prefix} - Elemento de texto XML',
                        'context': f'Texto: {match[:100]}...' if len(match) > 100 else f'Texto: {match}'
                    })
        
        # Buscar en CDATA sections que pueden contener URLs
        cdata_pattern = r'<!\[CDATA\[(.*?)\]\]>'
        cdata_matches = re.findall(cdata_pattern, xml_content, re.DOTALL)
        for cdata_content in cdata_matches:
            urls_in_cdata = self._find_urls_in_text(cdata_content)
            for url in urls_in_cdata:
                urls_found.append({
                    'url': url,
                    'location': f'{location_prefix} - CDATA',
                    'context': f'CDATA: {cdata_content[:100]}...' if len(cdata_content) > 100 else f'CDATA: {cdata_content}'
                })
        
        return urls_found
    
    def _find_urls_in_text(self, text):
        """Encontrar todas las URLs en un texto usando múltiples patrones"""
        if not text:
            return []
        
        urls_found = set()  # Usar set para evitar duplicados
        
        # Limpiar el texto de caracteres de control y espacios múltiples
        clean_text = re.sub(r'\s+', ' ', text.strip())
        
        # PASO 1: Buscar URLs completas primero (más específicas)
        # URLs con protocolo completo - estas tienen prioridad
        full_url_pattern = re.compile(r'https?://[^\s<>"\']+', re.IGNORECASE)
        full_urls = full_url_pattern.findall(clean_text)
        
        for url in full_urls:
            # Limpiar caracteres finales problemáticos
            cleaned_url = url.rstrip('.,;:)"\'>')
            if self._is_valid_url(cleaned_url):
                # Agregar protocolo si falta al final debido a limpieza
                if not cleaned_url.startswith(('http://', 'https://', 'ftp://')):
                    if '://' in cleaned_url:
                        cleaned_url = url  # Usar original si tiene protocolo
                urls_found.add(cleaned_url)
        
        # PASO 2: Buscar otros patrones solo si no encontramos la URL completa
        for pattern in self.url_patterns[1:]:  # Omitir el primer patrón ya usado
            matches = pattern.findall(clean_text)
            for match in matches:
                if isinstance(match, tuple):
                    # Para patrones que capturan grupos, tomar el grupo apropiado
                    url = match[1] if len(match) > 1 else match[0]
                else:
                    url = match
                
                if url:
                    # Limpiar la URL
                    url = url.strip().rstrip('.,;:)"\'>')
                    
                    # Verificar que no sea un fragmento de una URL ya encontrada
                    is_fragment = False
                    for existing_url in urls_found:
                        if url in existing_url and len(url) < len(existing_url):
                            is_fragment = True
                            break
                    
                    if not is_fragment and self._is_valid_url(url):
                        # Agregar protocolo si no lo tiene
                        if not url.startswith(('http://', 'https://', 'ftp://')):
                            if url.startswith('www.'):
                                url = 'http://' + url
                            elif '.' in url and not '@' in url:  # No es email
                                url = 'http://' + url
                        
                        urls_found.add(url)
        
        # PASO 3: Usar el patrón principal como respaldo SOLO para URLs no encontradas
        main_matches = self.url_pattern.findall(clean_text)
        for url in main_matches:
            url = url.strip().rstrip('.,;:)"\'>')
            
            # Verificar que no sea fragmento
            is_fragment = False
            for existing_url in urls_found:
                if url in existing_url and len(url) < len(existing_url):
                    is_fragment = True
                    break
            
            if not is_fragment and self._is_valid_url(url):
                # Agregar protocolo si no lo tiene
                if not url.startswith(('http://', 'https://', 'ftp://')):
                    if url.startswith('www.'):
                        url = 'http://' + url
                    elif '.' in url and not '@' in url:
                        url = 'http://' + url
                
                urls_found.add(url)
        
        # PASO 4: Post-procesamiento para eliminar fragmentos
        final_urls = []
        sorted_urls = sorted(urls_found, key=len, reverse=True)  # Más largas primero
        
        for url in sorted_urls:
            # Verificar que no sea un fragmento de una URL más larga ya incluida
            is_fragment = False
            for existing_url in final_urls:
                if url in existing_url and len(url) < len(existing_url):
                    is_fragment = True
                    break
            
            if not is_fragment:
                final_urls.append(url)
        
        return final_urls
    
    def _is_valid_url(self, url):
        """Validar si una URL es válida y no es un falso positivo"""
        if not url or len(url) < 4:
            return False
        
        # Limpiar la URL de caracteres extraños al final
        url = url.rstrip('.,;:)"\'>')
        
        # EXCLUIR COMPLETAMENTE METADATOS Y SCHEMAS - MUY RESTRICTIVO
        schema_patterns = [
            'http://schemas.openxmlformats.org',
            'http://schemas.microsoft.com',
            'http://www.w3.org',  # TODOS los w3.org
            'https://www.w3.org',  # TODOS los w3.org con https también
            'http://schemas.xmlsoap.org',
            'urn:schemas-microsoft-com',
            'xmlns:',
            'http://purl.org/dc/elements',
            'http://purl.org/dc/terms',
            'http://purl.org/dc/dcmitype',
            # SCHEMAS REPORTADOS POR EL USUARIO
            'http://customschemas.google.com',
            'http://customooxmlschemas.google.com',
            'http://ns.adobe.com',
            'urn:',
            'xmlns'
        ]
        
        # VERIFICAR SCHEMAS ESTRICTO
        for schema in schema_patterns:
            if url.startswith(schema):
                return False
        
        # EXCLUIR URLs que contengan caracteres de control o nulos
        if '\x00' in url or any(ord(c) < 32 for c in url):
            return False
        
        # EXCLUIR COMPLETAMENTE METADATOS DE ARCHIVOS PPT
        metadata_patterns = [
            'docProps/',
            '/docProps/',
            'core.xml',
            'app.xml',
            'custom.xml',
            '.xml',
            '.rels',
            '_rels',
            'metadata',
            'properties'
        ]
        
        url_lower = url.lower()
        for pattern in metadata_patterns:
            if pattern in url_lower:
                return False
        
        # Excluir falsos positivos comunes
        false_positives = [
            'example.com', 'test.com', 'localhost', '127.0.0.1',
            'your-domain.com', 'yoursite.com', 'website.com',
            'domain.com', 'site.com', 'company.com', 'sample.com',
            'placeholder.com', 'dummy.com', 'fake.com'
        ]
        
        for fp in false_positives:
            if fp in url_lower:
                return False
        
        # La URL debe tener al menos un punto para ser válida
        if '.' not in url:
            return False
        
        # SOLO PERMITIR URLs QUE SEAN CLARAMENTE VÁLIDAS Y REALES
        valid_patterns = [
            # URLs con protocolos de sitios reales
            r'^https?://(?:www\.)?(?:youtube|google|github|microsoft|amazon|facebook|twitter|linkedin|instagram|tiktok|vimeo|zoom|teams|stackoverflow|reddit|wikipedia|netflix|apple|adobe|oracle|ibm|salesforce|dropbox|slack|discord|whatsapp|telegram)\.com',
            # Dominios educativos
            r'^https?://[^/]+\.edu',
            # Dominios organizacionales
            r'^https?://[^/]+\.org',
            # Dominios gubernamentales
            r'^https?://[^/]+\.gov',
            # URLs que claramente son sitios web reales (con al menos 2 niveles de dominio)
            r'^https?://[a-zA-Z0-9-]+\.[a-zA-Z0-9-]+\.[a-zA-Z]{2,}',
            # IPs válidas
            r'^https?://(?:[0-9]{1,3}\.){3}[0-9]{1,3}'
        ]
        
        for pattern in valid_patterns:
            if re.search(pattern, url):
                return True
        
        return False
    
    def get_url_statistics(self, urls_list):
        """Obtener estadísticas detalladas de las URLs encontradas"""
        if not urls_list:
            return {}
        
        domains = [urlparse(url_info['url']).netloc for url_info in urls_list]
        locations = [url_info['location'] for url_info in urls_list]
        
        # Extraer números de diapositivas
        slide_numbers = []
        for url_info in urls_list:
            slide_match = re.search(r'Diapositiva (\d+)', url_info['location'])
            if slide_match:
                slide_numbers.append(int(slide_match.group(1)))
        
        return {
            'total_urls': len(urls_list),
            'unique_domains': len(set(domains)),
            'slides_with_urls': len(set(slide_numbers)),
            'most_common_domains': self._get_most_common_domains(domains),
            'by_location': self._group_by_location(urls_list),
            'by_slide': self._group_by_slide(urls_list)
        }
    
    def _get_most_common_domains(self, domains):
        """Obtener los dominios más comunes"""
        domain_count = {}
        for domain in domains:
            domain_count[domain] = domain_count.get(domain, 0) + 1
        
        return sorted(domain_count.items(), key=lambda x: x[1], reverse=True)[:10]
    
    def _group_by_location(self, urls_list):
        """Agrupar URLs por ubicación"""
        location_count = {}
        for url_info in urls_list:
            location = url_info['location']
            location_count[location] = location_count.get(location, 0) + 1
        
        return sorted(location_count.items(), key=lambda x: x[1], reverse=True)
    
    def _group_by_slide(self, urls_list):
        """Agrupar URLs por número de diapositiva"""
        slide_count = {}
        for url_info in urls_list:
            slide_match = re.search(r'Diapositiva (\d+)', url_info['location'])
            if slide_match:
                slide_num = int(slide_match.group(1))
                slide_count[slide_num] = slide_count.get(slide_num, 0) + 1
        
        return sorted(slide_count.items(), key=lambda x: x[0])
    
    def _extract_from_all_content_brute_force(self, zip_content):
        """BÚSQUEDA SELECTIVA en archivos relevantes (NO metadatos)"""
        urls_found = []
        
        try:
            with zipfile.ZipFile(io.BytesIO(zip_content), 'r') as zip_file:
                # Procesar SOLO archivos relevantes (NO docProps, NO _rels generales)
                for file_name in zip_file.namelist():
                    # EXCLUIR archivos de metadatos y propiedades
                    if any(excluded in file_name.lower() for excluded in [
                        'docprops/', 'docprops\\', 'core.xml', 'app.xml', 'custom.xml',
                        'metadata', 'properties', 'thumbnail'
                    ]):
                        continue
                    
                    # SOLO procesar archivos de slides y relaciones de slides
                    if not (file_name.startswith('ppt/slides/') or 
                           (file_name.endswith('.rels') and 'slides' in file_name)):
                        continue
                    
                    try:
                        # Leer contenido como texto
                        content = zip_file.read(file_name)
                        
                        # Intentar decodificar como texto
                        text_content = ""
                        for encoding in ['utf-8', 'utf-16', 'latin-1', 'cp1252']:
                            try:
                                text_content = content.decode(encoding)
                                break
                            except UnicodeDecodeError:
                                continue
                        
                        if text_content:
                            # Buscar URLs usando TODOS los patrones
                            urls_in_content = self._find_urls_in_text(text_content)
                            for url in urls_in_content:
                                # DOBLE VERIFICACIÓN de que sea válida
                                if self._is_valid_url(url):
                                    # Determinar si es un slide y cuál
                                    slide_num = 0
                                    if 'slide' in file_name:
                                        slide_match = re.search(r'slide(\d+)', file_name)
                                        slide_num = int(slide_match.group(1)) if slide_match else 0
                                    
                                    urls_found.append({
                                        'url': url,
                                        'location': f'Diapositiva {slide_num} - Búsqueda profunda' if slide_num > 0 else f'Archivo de slides - Búsqueda profunda',
                                        'context': f'Encontrado en búsqueda selectiva de {file_name}'
                                    })
                        
                    except Exception:
                        continue
                
                # NO buscar en archivos binarios para evitar metadatos
        
        except Exception as e:
            print(f"Error en búsqueda selectiva: {str(e)}")
        
        return urls_found
    
    def _deduplicate_urls_advanced(self, urls_found):
        """Deduplicación avanzada de URLs con múltiples criterios"""
        if not urls_found:
            return []
        
        # Paso 1: Crear un diccionario para agrupar por URL
        url_groups = {}
        
        for url_info in urls_found:
            url = url_info['url']
            
            # Normalizar la URL para comparación
            normalized_url = self._normalize_url_for_comparison(url)
            
            if normalized_url not in url_groups:
                url_groups[normalized_url] = []
            
            url_groups[normalized_url].append(url_info)
        
        # Paso 2: Para cada grupo de URLs similares, elegir la mejor
        unique_urls = []
        
        for normalized_url, group in url_groups.items():
            if len(group) == 1:
                # Solo una URL, mantenerla
                unique_urls.append(group[0])
            else:
                # Múltiples URLs, elegir la mejor
                best_url = self._choose_best_url_from_group(group)
                unique_urls.append(best_url)
        
        # Paso 3: Verificar fragmentos dentro de URLs más largas
        final_urls = self._remove_url_fragments(unique_urls)
        
        return final_urls
    
    def _normalize_url_for_comparison(self, url):
        """Normalizar URL para comparación de duplicados"""
        # Limpiar espacios y caracteres extraños
        url = url.strip().rstrip('.,;:)"\'>')
        
        # Convertir a minúsculas para comparación
        url = url.lower()
        
        # Quitar fragmentos finales comunes
        if url.endswith('/'):
            url = url[:-1]
        
        return url
    
    def _choose_best_url_from_group(self, group):
        """Elegir la mejor URL de un grupo de duplicados"""
        # Criterios de prioridad:
        # 1. URL más larga (más completa)
        # 2. URL con protocolo https sobre http
        # 3. URL con contexto más informativo
        
        # Ordenar por longitud (más larga primero)
        group_sorted = sorted(group, key=lambda x: len(x['url']), reverse=True)
        
        # Si hay URLs con https, priorizar esas
        https_urls = [url_info for url_info in group_sorted if url_info['url'].startswith('https://')]
        if https_urls:
            # Tomar la URL https más larga
            return https_urls[0]
        
        # Si no hay https, tomar la más larga
        return group_sorted[0]
    
    def _remove_url_fragments(self, urls_list):
        """Remover URLs que son fragmentos de otras URLs más largas"""
        if not urls_list:
            return []
        
        # Ordenar por longitud de URL (más largas primero)
        sorted_urls = sorted(urls_list, key=lambda x: len(x['url']), reverse=True)
        
        final_urls = []
        
        for current_url_info in sorted_urls:
            current_url = current_url_info['url'].lower()
            
            # Verificar si esta URL es un fragmento de alguna URL ya incluida
            is_fragment = False
            
            for existing_url_info in final_urls:
                existing_url = existing_url_info['url'].lower()
                
                # Si la URL actual está contenida en una URL más larga ya incluida
                if current_url in existing_url and len(current_url) < len(existing_url):
                    # Verificar que realmente sea un fragmento, no solo una coincidencia
                    if (current_url.startswith(existing_url[:len(current_url)]) or 
                        existing_url.startswith(current_url)):
                        is_fragment = True
                        break
            
            if not is_fragment:
                final_urls.append(current_url_info)
        
        return final_urls

# Función de utilidad para uso directo
def extract_urls_from_pptx(file_path_or_content):
    """
    Función auxiliar para extraer URLs de un archivo PPTX
    Mantener compatibilidad con código existente
    """
    extractor = PPTXURLExtractor()
    return extractor.extract_urls_from_file(file_path_or_content) 